"""Bullets: character bullets, auto-numbering schemes, bullet color/size,
and multi-level indentation — all raw pPr DrawingML (python-pptx has no
bullet API)."""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import blank_slide, grid_boxes, new_deck, ppr_append, save
from pptx.util import Pt

FONT = "Calibri"


def bullet_par(tf, text, first, level=0, bullets=()):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    for xml in bullets:
        ppr_append(p, xml)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.name = FONT
    return p


def main():
    prs = new_deck()
    slide = blank_slide(prs)
    cells = list(grid_boxes(4, cols=2))

    # character bullets (Wingdings and Arial)
    tb = slide.shapes.add_textbox(*cells[0])
    tf = tb.text_frame
    tf.word_wrap = True
    wingding = ('<a:buFont typeface="Wingdings"/>', '<a:buChar char="§"/>')
    dash = ('<a:buFont typeface="Arial"/>', '<a:buChar char="-"/>')
    bullet_par(tf, "Wingdings square bullet", True, bullets=wingding)
    bullet_par(tf, "Second bulleted line", False, bullets=wingding)
    bullet_par(tf, "Dash bullet line", False, bullets=dash)

    # auto-numbering schemes
    tb = slide.shapes.add_textbox(*cells[1])
    tf = tb.text_frame
    tf.word_wrap = True
    for i, scheme in enumerate(("arabicPeriod", "romanUcPeriod", "alphaLcParenR")):
        for j in range(2):
            bullet_par(tf, f"{scheme} item {j + 1}", i == 0 and j == 0,
                       bullets=(f'<a:buAutoNum type="{scheme}"/>',))

    # bullet color + size overrides
    tb = slide.shapes.add_textbox(*cells[2])
    tf = tb.text_frame
    tf.word_wrap = True
    styled = ('<a:buClr><a:srgbClr val="C00000"/></a:buClr>',
              '<a:buSzPct val="150000"/>',
              '<a:buFont typeface="Arial"/>', '<a:buChar char="●"/>')
    bullet_par(tf, "Red oversized bullet", True, bullets=styled)
    bullet_par(tf, "Another red bullet line", False, bullets=styled)

    # multi-level indentation
    tb = slide.shapes.add_textbox(*cells[3])
    tf = tb.text_frame
    tf.word_wrap = True
    dot = ('<a:buFont typeface="Arial"/>', '<a:buChar char="•"/>')
    for level in range(4):
        bullet_par(tf, f"Indent level {level}", level == 0, level=level,
                   bullets=dot)

    save(prs)


if __name__ == "__main__":
    main()
