"""Text alignment, wrapping, vertical anchoring, and overflow.

Boxes get a thin gray outline so their bounds are visible in both renders —
overflow past the outline is the text-clipping parity signal.
"""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import blank_slide, grid_boxes, new_deck, save
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

LOREM = ("The quick brown fox jumps over the lazy dog while seventy-seven "
         "wizards briskly mix quartz vials of juniper extract.")
FONT = "Calibri"


def outline(box):
    box.line.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
    box.line.width = Pt(0.75)


def textbox(slide, cell, text, align=None, anchor=None, wrap=True, size=14):
    tb = slide.shapes.add_textbox(*cell)
    tf = tb.text_frame
    tf.word_wrap = wrap
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = FONT
    if align is not None:
        p.alignment = align
    outline(tb)
    return tb


def main():
    prs = new_deck()

    slide = blank_slide(prs)
    cells = list(grid_boxes(8, cols=2))
    textbox(slide, cells[0], LOREM, align=PP_ALIGN.LEFT)
    textbox(slide, cells[1], LOREM, align=PP_ALIGN.CENTER)
    textbox(slide, cells[2], LOREM, align=PP_ALIGN.RIGHT)
    textbox(slide, cells[3], LOREM, align=PP_ALIGN.JUSTIFY)
    textbox(slide, cells[4], LOREM, anchor=MSO_ANCHOR.TOP, size=11)
    textbox(slide, cells[5], LOREM, anchor=MSO_ANCHOR.MIDDLE, size=11)
    textbox(slide, cells[6], LOREM, anchor=MSO_ANCHOR.BOTTOM, size=11)
    # overflow: text far too large for the box — PowerPoint keeps it inside
    # the box footprint only via clipping in export
    textbox(slide, cells[7], LOREM + " " + LOREM, size=24)

    slide = blank_slide(prs)
    cells = list(grid_boxes(4, cols=2))
    textbox(slide, cells[0], "No wrap: " + LOREM, wrap=False)
    textbox(slide, cells[1], "Word wrap on.\n" + LOREM)
    # multi-paragraph with a long unbreakable token near the wrap boundary
    tb = textbox(slide, cells[2], "Boundary word test:")
    p2 = tb.text_frame.add_paragraph()
    r = p2.add_run()
    r.text = "Antidisestablishmentarianism supercalifragilisticexpialidocious"
    r.font.size = Pt(16)
    r.font.name = FONT
    textbox(slide, cells[3], LOREM, align=PP_ALIGN.JUSTIFY, size=18)

    save(prs)


if __name__ == "__main__":
    main()
