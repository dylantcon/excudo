"""Gradient fills: linear at several angles, multi-stop, radial/rect/shape
path gradients, and stop-level alpha/modifiers.

Slide 1 uses python-pptx's gradient API (linear, two stops, angles).
Slide 2 hand-authors gradFill XML for everything the API can't express.
"""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import blank_slide, grid_boxes, new_deck, save, set_raw_fill
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ANGLES = [0, 45, 90, 135, 270]

RAW_GRADIENTS = [
    # 3-stop linear with a transparent middle stop
    ('<a:gradFill rotWithShape="1"><a:gsLst>'
     '<a:gs pos="0"><a:srgbClr val="C00000"/></a:gs>'
     '<a:gs pos="50000"><a:srgbClr val="FFC000"><a:alpha val="40000"/></a:srgbClr></a:gs>'
     '<a:gs pos="100000"><a:srgbClr val="1F4E79"/></a:gs>'
     '</a:gsLst><a:lin ang="2700000" scaled="1"/></a:gradFill>'),
    # centered radial
    ('<a:gradFill><a:gsLst>'
     '<a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>'
     '<a:gs pos="100000"><a:srgbClr val="2E74B5"/></a:gs>'
     '</a:gsLst><a:path path="circle">'
     '<a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path></a:gradFill>'),
    # corner-focused radial
    ('<a:gradFill><a:gsLst>'
     '<a:gs pos="0"><a:srgbClr val="FFD966"/></a:gs>'
     '<a:gs pos="100000"><a:srgbClr val="843C0C"/></a:gs>'
     '</a:gsLst><a:path path="circle">'
     '<a:fillToRect l="100000" t="100000"/></a:path></a:gradFill>'),
    # rect path
    ('<a:gradFill><a:gsLst>'
     '<a:gs pos="0"><a:srgbClr val="E2EFDA"/></a:gs>'
     '<a:gs pos="100000"><a:srgbClr val="375623"/></a:gs>'
     '</a:gsLst><a:path path="rect">'
     '<a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path></a:gradFill>'),
    # shape path on a rounded rectangle
    ('<a:gradFill><a:gsLst>'
     '<a:gs pos="0"><a:srgbClr val="FBE5D6"/></a:gs>'
     '<a:gs pos="100000"><a:srgbClr val="C55A11"/></a:gs>'
     '</a:gsLst><a:path path="shape">'
     '<a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path></a:gradFill>'),
    # scheme-color stops with modifiers
    ('<a:gradFill rotWithShape="1"><a:gsLst>'
     '<a:gs pos="0"><a:schemeClr val="accent1"><a:lumMod val="40000"/>'
     '<a:lumOff val="60000"/></a:schemeClr></a:gs>'
     '<a:gs pos="100000"><a:schemeClr val="accent1"><a:shade val="50000"/>'
     '</a:schemeClr></a:gs>'
     '</a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill>'),
]


def main():
    prs = new_deck()

    slide = blank_slide(prs)
    cells = list(grid_boxes(6, cols=3))
    for angle, box in zip(ANGLES, cells):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *box)
        shape.fill.gradient()
        shape.fill.gradient_stops[0].color.rgb = RGBColor(0xC0, 0x00, 0x00)
        shape.fill.gradient_stops[1].color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        shape.fill.gradient_angle = angle
        shape.line.fill.background()

    slide = blank_slide(prs)
    cells = list(grid_boxes(6, cols=3))
    for i, (grad, box) in enumerate(zip(RAW_GRADIENTS, cells)):
        preset = MSO_SHAPE.ROUNDED_RECTANGLE if i == 4 else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(preset, *box)
        set_raw_fill(shape, grad)
        shape.line.fill.background()

    save(prs)


if __name__ == "__main__":
    main()
