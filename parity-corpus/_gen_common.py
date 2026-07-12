"""Shared helpers for parity-corpus deck generators.

Every category's gen_deck.py runs with cwd set to its own directory and
writes ./deck.pptx. Decks are 16:9 (1280x720 px at 96 DPI), built on the
blank layout, using Calibri so ground truth and render share a font.

python-pptx covers most authoring; where its API stops (gradient path
types, alpha, arrowheads, bullets, effects, picture fills) the helpers
below splice raw DrawingML into the shape XML.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu

SLIDE_W_EMU = 12192000  # 1280 px
SLIDE_H_EMU = 6858000   # 720 px
MARGIN_EMU = 457200     # 0.5 in

# spPr fill elements, in schema order (only one may be present at a time)
_FILL_TAGS = ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
              "a:pattFill", "a:grpFill")


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank


def save(prs: Presentation) -> None:
    out = Path.cwd() / "deck.pptx"
    prs.save(str(out))
    print(f"wrote {out}")


def grid_boxes(count_per_slide: int = 24, cols: int = 6):
    """(left, top, width, height) EMU boxes for a cols-wide grid, laid out
    inside the slide margins. Yields exactly count_per_slide boxes."""
    rows = (count_per_slide + cols - 1) // cols
    cell_w = (SLIDE_W_EMU - 2 * MARGIN_EMU) // cols
    cell_h = (SLIDE_H_EMU - 2 * MARGIN_EMU) // rows
    pad_w, pad_h = cell_w // 8, cell_h // 8
    for i in range(count_per_slide):
        r, c = divmod(i, cols)
        yield (MARGIN_EMU + c * cell_w + pad_w,
               MARGIN_EMU + r * cell_h + pad_h,
               cell_w - 2 * pad_w,
               cell_h - 2 * pad_h)


def chunked(seq, n):
    for i in range(0, len(seq), n):
        yield seq[i:i + n]


# ---------- raw DrawingML splicing ----------

def _spPr(shape):
    return shape._element.spPr


def set_raw_fill(shape, fill_xml: str) -> None:
    """Replace the shape's fill with raw DrawingML (a:solidFill/a:gradFill/
    a:pattFill/a:blipFill... markup, 'a' namespace assumed)."""
    spPr = _spPr(shape)
    for tag in _FILL_TAGS:
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    fill = parse_xml(_with_ns(fill_xml))
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(fill)
    else:
        spPr.append(fill)


def _with_ns(xml: str) -> str:
    """Inject the drawingml namespace declarations into the root element,
    handling both <a:x ...>...</a:x> and self-closing <a:x .../> roots."""
    xml = xml.strip()
    i = xml.index(">")
    if xml[i - 1] == "/":
        return f"{xml[:i - 1]} {nsdecls('a', 'r')}/{xml[i:]}"
    return f"{xml[:i]} {nsdecls('a', 'r')}{xml[i:]}"


def raw_xml(xml: str):
    """Parse a DrawingML fragment ('a'/'r' namespaces assumed) into an lxml
    element ready to append into shape XML."""
    return parse_xml(_with_ns(xml))


def append_effects(shape, effect_lst_xml: str) -> None:
    """Append an a:effectLst (or a:effectDag) to the shape's spPr."""
    _spPr(shape).append(parse_xml(_with_ns(effect_lst_xml)))


def set_line_ends(connector, head: str | None = None, tail: str | None = None,
                  size: str = "med") -> None:
    """Set a:headEnd / a:tailEnd on a connector or shape outline.
    `head`/`tail` are OOXML types: triangle, stealth, diamond, oval, arrow."""
    ln = connector.line._get_or_add_ln()
    for name, kind in (("a:headEnd", head), ("a:tailEnd", tail)):
        existing = ln.find(qn(name))
        if existing is not None:
            ln.remove(existing)
        if kind:
            ln.append(parse_xml(
                _with_ns(f'<{name} type="{kind}" w="{size}" len="{size}"/>')))


def rpr_attrs(run, **attrs) -> None:
    """Set raw attributes on the run's rPr (e.g. spc='300', cap='small')."""
    rPr = run._r.get_or_add_rPr()
    for k, v in attrs.items():
        rPr.set(k, str(v))


def ppr_append(paragraph, xml: str) -> None:
    """Append a raw child element (buChar, buClr, ...) to the paragraph pPr."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.append(parse_xml(_with_ns(xml)))


def body_autofit(text_frame, font_scale: int, lnspc_reduction: int = 0) -> None:
    """Author a normAutofit with PowerPoint-style precomputed percentages
    (thousandths: 62500 = 62.5%)."""
    bodyPr = text_frame._txBody.bodyPr
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)
    xml = f'<a:normAutofit fontScale="{font_scale}"'
    if lnspc_reduction:
        xml += f' lnSpcReduction="{lnspc_reduction}"'
    xml += "/>"
    bodyPr.append(parse_xml(_with_ns(xml)))


def picture_fill(shape, image_path: Path, mode: str = "stretch") -> None:
    """Give the shape an a:blipFill referencing an embedded image part.
    mode: 'stretch' or 'tile'."""
    slide_part = shape.part
    image_part, rId = slide_part.get_or_add_image_part(str(image_path))
    inner = ('<a:stretch><a:fillRect/></a:stretch>' if mode == "stretch"
             else '<a:tile tx="0" ty="0" sx="50000" sy="50000" flip="none" algn="tl"/>')
    set_raw_fill(shape, f'<a:blipFill><a:blip r:embed="{rId}"/>{inner}</a:blipFill>')


def preset_shape_partition() -> dict:
    """Partition every MSO_SHAPE member into the six preset-shape corpus
    categories. Driven off the enum itself so the corpus always covers the
    full preset vocabulary python-pptx can author — no hand-listed names to
    rot."""
    from pptx.enum.shapes import MSO_SHAPE

    cats = {"preset-shapes-actionbuttons": [], "preset-shapes-flowchart": [],
            "preset-shapes-stars-banners": [], "preset-shapes-callouts": [],
            "preset-shapes-arrows": [], "preset-shapes-basic": []}
    for name, member in sorted(MSO_SHAPE.__members__.items()):
        if member.value < 1:  # MIXED and friends — not real presets
            continue
        if name.startswith("ACTION_BUTTON"):
            cats["preset-shapes-actionbuttons"].append(member)
        elif name.startswith("FLOWCHART"):
            cats["preset-shapes-flowchart"].append(member)
        elif any(k in name for k in ("STAR", "RIBBON", "SCROLL", "WAVE", "EXPLOSION")):
            cats["preset-shapes-stars-banners"].append(member)
        elif "CALLOUT" in name:
            cats["preset-shapes-callouts"].append(member)
        elif "ARROW" in name:
            cats["preset-shapes-arrows"].append(member)
        else:
            cats["preset-shapes-basic"].append(member)
    return cats


def die(msg: str) -> None:
    print(f"gen_deck failed: {msg}", file=sys.stderr)
    sys.exit(1)
