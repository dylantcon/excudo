"""Line styles: dash presets, widths, caps, joins, compound lines.

Slide 1: dash styles x widths on connector lines (python-pptx API).
Slide 2: raw a:ln attributes — cap (rnd/sq/flat), join (round/bevel/miter),
compound (dbl/thickThin) — on thick-outlined shapes where the difference
is visible at corners and line ends.
"""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import blank_slide, grid_boxes, new_deck, raw_xml, save
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt

# every real dash preset the enum offers (sentinels have negative values)
DASHES = [m for _, m in sorted(MSO_LINE_DASH_STYLE.__members__.items())
          if m.value >= 1]
WIDTHS = [Pt(1), Pt(3), Pt(6)]


def _ln(shape):
    return shape.line._get_or_add_ln()


def main():
    prs = new_deck()

    slide = blank_slide(prs)
    cells = list(grid_boxes(len(DASHES) * len(WIDTHS), cols=3))
    i = 0
    for width in WIDTHS:
        for dash in DASHES:
            left, top, w, h = cells[i]
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, left, top + h // 2, left + w, top + h // 2)
            conn.line.width = width
            conn.line.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            conn.line.dash_style = dash
            i += 1

    slide = blank_slide(prs)
    cells = list(grid_boxes(8, cols=4))
    # caps on fat open lines
    for cap, box in zip(("rnd", "sq", "flat"), cells):
        left, top, w, h = box
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, left, top + h // 2, left + w, top + h // 2)
        conn.line.width = Pt(14)
        conn.line.color.rgb = RGBColor(0xC0, 0x00, 0x00)
        _ln(conn).set("cap", cap)
    # joins on thick-outlined triangles
    joins = ["<a:round/>", "<a:bevel/>", '<a:miter lim="800000"/>']
    for join, box in zip(joins, cells[3:]):
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, *box)
        shape.fill.background()
        shape.line.width = Pt(10)
        shape.line.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
        _ln(shape).append(raw_xml(join))
    # compound lines
    for cmpd, box in zip(("dbl", "thickThin"), cells[6:]):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *box)
        shape.fill.background()
        shape.line.width = Pt(8)
        shape.line.color.rgb = RGBColor(0x53, 0x82, 0x35)
        _ln(shape).set("cmpd", cmpd)

    save(prs)


if __name__ == "__main__":
    main()
