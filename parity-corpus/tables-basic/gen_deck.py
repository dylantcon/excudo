"""Basic table: header row, banded rows (default table style), text in
cells, mixed column widths."""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import MARGIN_EMU, SLIDE_H_EMU, SLIDE_W_EMU, blank_slide, new_deck, save
from pptx.util import Emu, Pt

FONT = "Calibri"
HEADERS = ["Region", "Q1", "Q2", "Q3"]
ROWS = [
    ["North", "104", "121", "97"],
    ["South", "88", "95", "110"],
    ["East", "132", "128", "140"],
    ["West", "76", "83", "91"],
]


def main():
    prs = new_deck()
    slide = blank_slide(prs)
    left = MARGIN_EMU
    top = MARGIN_EMU
    width = SLIDE_W_EMU - 2 * MARGIN_EMU
    height = (SLIDE_H_EMU - 2 * MARGIN_EMU) * 2 // 3
    table = slide.shapes.add_table(len(ROWS) + 1, len(HEADERS),
                                   left, top, width, height).table
    table.columns[0].width = Emu(width // 2)
    for c, text in enumerate(HEADERS):
        cell = table.cell(0, c)
        cell.text = text
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].runs[0].font.name = FONT
    for r, row in enumerate(ROWS, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].runs[0].font.name = FONT
    save(prs)


if __name__ == "__main__":
    main()
