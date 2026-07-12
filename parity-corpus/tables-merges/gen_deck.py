"""Table cell merges: horizontal span, vertical span, and a block merge,
plus per-cell fill overrides."""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import MARGIN_EMU, SLIDE_H_EMU, SLIDE_W_EMU, blank_slide, new_deck, save
from pptx.dml.color import RGBColor
from pptx.util import Pt

FONT = "Calibri"


def main():
    prs = new_deck()
    slide = blank_slide(prs)
    width = SLIDE_W_EMU - 2 * MARGIN_EMU
    height = (SLIDE_H_EMU - 2 * MARGIN_EMU) * 3 // 4
    table = slide.shapes.add_table(5, 5, MARGIN_EMU, MARGIN_EMU,
                                   width, height).table

    for r in range(5):
        for c in range(5):
            cell = table.cell(r, c)
            cell.text = f"r{r}c{c}"
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].runs[0].font.name = FONT

    # horizontal merge across the top row
    table.cell(0, 0).merge(table.cell(0, 4))
    table.cell(0, 0).text = "Merged header row"
    # vertical merge down the first column
    table.cell(1, 0).merge(table.cell(4, 0))
    table.cell(1, 0).text = "Tall"
    # 2x2 block merge in the middle
    table.cell(2, 2).merge(table.cell(3, 3))
    table.cell(2, 2).text = "2x2 block"
    # per-cell fill override
    override = table.cell(1, 4)
    override.fill.solid()
    override.fill.fore_color.rgb = RGBColor(0xC0, 0x00, 0x00)

    for cell in (table.cell(0, 0), table.cell(1, 0), table.cell(2, 2)):
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(14)
        run.font.name = FONT

    save(prs)


if __name__ == "__main__":
    main()
