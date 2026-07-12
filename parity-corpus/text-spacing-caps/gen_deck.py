"""Character-level text properties: tracking (spc), small/all caps,
underline, strikethrough, super/subscript, and line spacing."""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import blank_slide, grid_boxes, new_deck, rpr_attrs, save
from pptx.dml.color import RGBColor
from pptx.util import Pt

FONT = "Calibri"

CASES = [
    ("Normal tracking baseline text", {}),
    ("Expanded tracking +3pt", {"spc": 300}),
    ("Tight tracking -1pt", {"spc": -100}),
    ("Small Caps Sample Text", {"cap": "small"}),
    ("all caps sample text", {"cap": "all"}),
    ("Underlined sample text", {"u": "sng"}),
    ("Struck-through sample text", {"strike": "sngStrike"}),
    ("Superscript sample", {"baseline": 30000}),
    ("Subscript sample", {"baseline": -25000}),
]

LINE_SPACINGS = [1.0, 1.5, 2.0]


def main():
    prs = new_deck()

    slide = blank_slide(prs)
    cells = list(grid_boxes(9, cols=3))
    for (text, attrs), cell in zip(CASES, cells):
        tb = slide.shapes.add_textbox(*cell)
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.name = FONT
        if attrs:
            rpr_attrs(run, **attrs)

    slide = blank_slide(prs)
    cells = list(grid_boxes(3, cols=3))
    for spacing, cell in zip(LINE_SPACINGS, cells):
        tb = slide.shapes.add_textbox(*cell)
        tf = tb.text_frame
        tf.word_wrap = True
        for i in range(3):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = spacing
            run = p.add_run()
            run.text = f"Line spacing {spacing:g} paragraph {i + 1} wraps here"
            run.font.size = Pt(14)
            run.font.name = FONT

    save(prs)


if __name__ == "__main__":
    main()
