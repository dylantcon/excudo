"""Theme-color fills x brightness, plus the raw color-modifier vocabulary.

Slide 1: python-pptx theme fills (accent1-6 x brightness -0.5..+0.5), which
PowerPoint stores as lumMod/lumOff on schemeClr — the single most common
color pattern in real themed decks.
Slide 2: hand-authored modifier matrix (lumMod, lumOff, shade, tint, satMod,
alpha) so each modifier is pinned individually.
"""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import blank_slide, grid_boxes, new_deck, save, set_raw_fill
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE

ACCENTS = [MSO_THEME_COLOR.ACCENT_1, MSO_THEME_COLOR.ACCENT_2,
           MSO_THEME_COLOR.ACCENT_3, MSO_THEME_COLOR.ACCENT_4,
           MSO_THEME_COLOR.ACCENT_5, MSO_THEME_COLOR.ACCENT_6]
BRIGHTNESS = [-0.5, -0.25, 0.0, 0.25, 0.5]

MODIFIERS = [
    '<a:lumMod val="75000"/>',
    '<a:lumMod val="50000"/>',
    '<a:lumMod val="60000"/><a:lumOff val="40000"/>',
    '<a:lumMod val="40000"/><a:lumOff val="60000"/>',
    '<a:shade val="50000"/>',
    '<a:tint val="50000"/>',
    '<a:satMod val="200000"/>',
    '<a:satMod val="50000"/>',
    '<a:alpha val="50000"/>',
    '<a:alpha val="25000"/>',
    '<a:shade val="75000"/><a:satMod val="150000"/>',
    '<a:tint val="75000"/><a:alpha val="60000"/>',
]
MOD_SCHEMES = ["accent1", "accent2"]


def main():
    prs = new_deck()

    slide = blank_slide(prs)
    cells = list(grid_boxes(30, cols=6))
    for i, brightness in enumerate(BRIGHTNESS):
        for j, accent in enumerate(ACCENTS):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *cells[i * 6 + j])
            shape.fill.solid()
            shape.fill.fore_color.theme_color = accent
            if brightness:
                shape.fill.fore_color.brightness = brightness
            shape.line.fill.background()

    slide = blank_slide(prs)
    cells = list(grid_boxes(24, cols=6))
    idx = 0
    for scheme in MOD_SCHEMES:
        for mods in MODIFIERS:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *cells[idx])
            set_raw_fill(shape, f'<a:solidFill><a:schemeClr val="{scheme}">'
                                f'{mods}</a:schemeClr></a:solidFill>')
            shape.line.fill.background()
            idx += 1

    save(prs)


if __name__ == "__main__":
    main()
