"""Charts (permanently expected-fail: real chart rendering is out of scope;
Excudo renders a labeled placeholder frame instead). Kept in the corpus so
the dashboard shows the divergence honestly."""
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent))

from _gen_common import blank_slide, grid_boxes, new_deck, save
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def main():
    prs = new_deck()
    slide = blank_slide(prs)
    cells = list(grid_boxes(2, cols=2))

    data = CategoryChartData()
    data.categories = ["North", "South", "East", "West"]
    data.add_series("Q1", (104, 88, 132, 76))
    data.add_series("Q2", (121, 95, 128, 83))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *cells[0], data)

    pie = CategoryChartData()
    pie.categories = ["A", "B", "C"]
    pie.add_series("Share", (45, 30, 25))
    slide.shapes.add_chart(XL_CHART_TYPE.PIE, *cells[1], pie)

    save(prs)


if __name__ == "__main__":
    main()
